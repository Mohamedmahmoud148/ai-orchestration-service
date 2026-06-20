"""
react_agent.py — ReAct Agent (gpt-4o-mini + native function calling)

Loop (max 3 iterations):
  Think → call tool(s) in parallel → observe result → think → ... → final answer

Performance improvements (Wave 3):
  - MAX_ITERATIONS reduced 6 → 3   (95% of queries finish in 1-2)
  - Context fast-path               (identity/profile answered without any LLM call)
  - Parallel tool calls             (asyncio.gather instead of sequential loop)
  - stream_run() async generator    (final answer streamed token-by-token)
"""
from __future__ import annotations

import asyncio
import json
import time
from typing import Any, AsyncGenerator, Dict, List, Optional, TYPE_CHECKING

from app.core.api_discovery import get_allowed_endpoints_schema, validate_endpoint
from app.core.logging import logger
from app.core.response_guard import validate as guard_response, check_user_input
from app.prompts import load_prompt

if TYPE_CHECKING:
    from app.agents.execution_context import ExecutionContext

_MODEL = "openai/gpt-4o-mini"
_MAX_ITERATIONS = 4   # 4 allows: think → tool1 → tool2 → final answer for complex academic queries

# Follow-up phrases that indicate the user wants to act on the previously active document
_FOLLOWUP_AR = (
    "لخصه", "لخصلي", "لخص", "اقراه", "اقرا", "اقرأه", "اقرأ", "اشرحه", "اشرح الملف",
    "اشرح المحتوى", "اشرح المحتوي", "ماذا يحتوي", "ما يحتوي", "اعمل quiz",
    "اعمل امتحان", "استخرج العناوين", "ملخص", "فيه ايه", "فيه إيه",
    "محتواه", "محتواها", "عاوز تلخيص", "عايز تلخيص",
)
_FOLLOWUP_EN = (
    "summarize it", "summarize this", "read it", "read this", "explain it",
    "explain this", "make a quiz", "generate exam", "generate a quiz",
    "list headings", "what's inside", "what is inside",
)

# ── Tool definitions ──────────────────────────────────────────────────────────

_TOOL_CALL_API: dict = {
    "type": "function",
    "function": {
        "name": "call_backend_api",
        "description": (
            "Call the university management system backend API to fetch or submit real data. "
            "Use this for ANY question that needs actual data from the system: students, grades, "
            "schedules, complaints, enrollments, analytics, departments, batches, exams, etc. "
            "You can call it multiple times — each call returns fresh data. "
            "NEVER guess or fabricate data; always fetch it."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "method": {
                    "type": "string",
                    "enum": ["GET", "POST"],
                    "description": "HTTP method",
                },
                "path": {
                    "type": "string",
                    "description": (
                        "Full API path, e.g. /api/Students or /api/Students/01KS3XG... "
                        "Substitute real IDs from the user context or previous tool results."
                    ),
                },
                "query_params": {
                    "type": "object",
                    "description": (
                        "Query string key-value pairs, "
                        "e.g. {\"page\": 1, \"size\": 20, \"userId\": \"01KS...\"}"
                    ),
                    "additionalProperties": True,
                },
                "body": {
                    "type": "object",
                    "description": "Request body for POST requests.",
                    "additionalProperties": True,
                },
            },
            "required": ["method", "path"],
        },
    },
}

_TOOL_REGULATION: dict = {
    "type": "function",
    "function": {
        "name": "read_regulation_pdf",
        "description": (
            "Read and search the official academic regulation / student handbook PDF. "
            "Use this when the user asks about: subjects per year, curriculum, credit hours, "
            "graduation requirements, study plan, academic policies, or anything from the دليل الطالب."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "What the user wants to know from the regulation document.",
                }
            },
            "required": ["query"],
        },
    },
}

_TOOL_GENERATE_EXAM: dict = {
    "type": "function",
    "function": {
        "name": "generate_exam",
        "description": (
            "Generate an AI-powered exam. Use ONLY when a doctor or admin explicitly asks "
            "to create/generate an exam. Requires: subject name, question count, difficulty, exam type."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "subject": {"type": "string", "description": "Subject/course name"},
                "question_count": {"type": "integer", "description": "Number of questions (5-50)"},
                "difficulty": {
                    "type": "string",
                    "enum": ["easy", "medium", "hard"],
                    "description": "Exam difficulty level",
                },
                "exam_type": {
                    "type": "string",
                    "enum": ["mcq", "truefalse", "essay", "mixed"],
                    "description": "Type of questions",
                },
                "topics": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Optional list of specific topics to focus on",
                },
                "subject_offering_id": {
                    "type": "string",
                    "description": "SubjectOffering ID if known (from context or previous API call)",
                },
            },
            "required": ["subject", "question_count", "difficulty", "exam_type"],
        },
    },
}

_TOOL_ACADEMIC_ANALYSIS: dict = {
    "type": "function",
    "function": {
        "name": "analyze_academic_profile",
        "description": (
            "Deep academic advisor analysis for a student. "
            "Use this when: student asks about GPA, academic risk, graduation plan, "
            "failed subjects, credit hours progress, recommended next subjects, "
            "or asks 'وضعي الأكاديمي إيه؟', 'هل هتخرج؟', 'أنا في خطر؟', "
            "'ايه المواد الباقية؟', 'كيف أحسن معدلي؟'. "
            "Fetches roadmap + grades + regulation passages and returns a deep analysis."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "focus": {
                    "type": "string",
                    "description": (
                        "What the student specifically wants to know. "
                        "E.g.: 'graduation_readiness', 'gpa_improvement', "
                        "'failed_subjects', 'next_semester_plan', 'general'"
                    ),
                    "enum": [
                        "graduation_readiness",
                        "gpa_improvement",
                        "failed_subjects",
                        "next_semester_plan",
                        "academic_risk",
                        "general",
                    ],
                },
                "question": {
                    "type": "string",
                    "description": "The student's specific question to focus the analysis on.",
                },
            },
            "required": ["focus"],
        },
    },
}

_TOOL_READ_MATERIAL: dict = {
    "type": "function",
    "function": {
        "name": "read_material_pdf",
        "description": (
            "Download and read the CONTENT INSIDE a specific PDF file whose URL is already known. "
            "ONLY use this when: (1) you already have the file_url from a previous step or the user's message, "
            "AND (2) the user explicitly asks to read/summarize/explain the file content — "
            "e.g.: 'لخصلي الملف', 'اقرا المحاضرة دي', 'ايه اللي جوا الملف ده', 'summarize this PDF'. "
            "NEVER use to list available materials or check what files exist — use call_backend_api for that. "
            "NEVER use for regulations/student handbook — use read_regulation_pdf for that. "
            "If file_url is unknown, first call the backend API to get it, then call this tool."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "file_url": {
                    "type": "string",
                    "description": "The direct URL of the PDF file to read. MUST be known before calling this tool.",
                },
                "task": {
                    "type": "string",
                    "description": "What to do: 'summarize' | 'list_headings' | 'explain' | 'read'",
                    "enum": ["summarize", "list_headings", "explain", "read"],
                },
                "question": {
                    "type": "string",
                    "description": "Specific question to answer from the file content (optional).",
                },
            },
            "required": ["file_url", "task"],
        },
    },
}

_ALL_TOOLS = [_TOOL_CALL_API, _TOOL_REGULATION, _TOOL_READ_MATERIAL, _TOOL_GENERATE_EXAM, _TOOL_ACADEMIC_ANALYSIS]


# ── Context fast-path ─────────────────────────────────────────────────────────

_IDENTITY_KEYWORDS = (
    "اسمي", "اسمك", "انا مين", "أنا مين", "مين انا", "مين أنا",
    "من انا", "من أنا", "who am i", "my name", "what is my name",
)

def _normalize_academic_context(ctx: dict) -> dict:
    """
    The .NET backend serializes with JsonNamingPolicy.SnakeCaseLower.
    E.g. batchId → batch_id.  ReactAgent expects camelCase keys.
    This function ensures BOTH forms exist so lookups always succeed.
    """
    if not ctx:
        return ctx
    mappings = [
        ("batch_id",       "batchId"),
        ("student_id",     "studentId"),
        ("user_id",        "userId"),
        ("doctor_id",      "doctorId"),
        ("department_id",  "departmentId"),
        ("college_id",     "collegeId"),
        ("group_id",       "groupId"),
        ("semester_id",    "semesterId"),
        ("offering_id",    "offeringId"),
        ("subject_offering_id", "subjectOfferingId"),
        ("batch_name",     "batchName"),
        ("student_name",   "studentName"),
        ("full_name",      "fullName"),
        ("doctor_name",    "doctorName"),
        ("department_name","departmentName"),
        ("college_name",   "collegeName"),
    ]
    for snake, camel in mappings:
        # If snake_case exists but camelCase doesn't → add camelCase alias
        if snake in ctx and camel not in ctx:
            ctx[camel] = ctx[snake]
        # If camelCase exists but snake_case doesn't → add snake_case alias
        elif camel in ctx and snake not in ctx:
            ctx[snake] = ctx[camel]
    return ctx


def _try_answer_from_context(context: "ExecutionContext") -> Optional[str]:
    """
    Answer simple identity/profile questions directly from academic_context
    without any LLM call — zero latency for the most common queries.
    """
    ctx = context.academic_context or {}
    msg = context.message.lower().strip()

    if any(kw in msg for kw in _IDENTITY_KEYWORDS):
        name = (
            ctx.get("fullName") or ctx.get("studentName")
            or ctx.get("name") or ctx.get("doctorName")
        )
        if name:
            return f"اسمك هو **{name}** 😊"

    return None


# ── System prompt builder ─────────────────────────────────────────────────────

def _load_role_persona(role: str) -> str:
    """Load role-specific personality from prompts/role_{role}.md. Silent fallback."""
    role_clean = (role or "").lower().strip()
    known = {"student", "doctor", "admin", "superadmin"}
    name = f"role_{role_clean}" if role_clean in known else "role_student"
    try:
        return load_prompt(name)
    except Exception:
        return ""


def _build_active_doc_section(active_doc: dict) -> str:
    """Build a system prompt section for the currently active document."""
    if not active_doc:
        return ""
    doc_type = active_doc.get("document_type", "material")
    title = active_doc.get("title", "")
    subject = active_doc.get("subject_name", "")
    file_url = active_doc.get("file_url", "")

    if doc_type == "material" and file_url:
        subject_line = f"\n- المادة: {subject}" if subject else ""
        return (
            f"\n## الوثيقة النشطة حالياً\n"
            f"- النوع: ملف مادة دراسية (material)\n"
            f"- العنوان: {title}{subject_line}\n"
            f"- الرابط: {file_url}\n"
            f"⚠️ **قاعدة ملزمة**: إذا قال المستخدم 'لخصه' أو 'اقراه' أو 'اشرحه' أو 'اشرح الملف' "
            f"أو 'اشرح المحتوى' أو 'ماذا يحتوي' أو 'اعمل quiz' أو 'اعمل امتحان' أو 'استخرج العناوين' "
            f"أو 'ملخص' أو 'summarize' أو 'read it' أو 'explain it' أو 'make a quiz' "
            f"→ استخدم `read_material_pdf` مع الرابط أعلاه مباشرةً. "
            f"لا تستخدم `read_regulation_pdf` أبداً لهذه الطلبات.\n"
        )
    elif doc_type == "regulation":
        return (
            f"\n## الوثيقة النشطة حالياً\n"
            f"- النوع: لائحة أكاديمية (regulation)\n"
            f"- العنوان: {title}\n"
            f"إذا طلب المستخدم مزيداً من التفاصيل → استخدم `read_regulation_pdf`.\n"
        )
    return ""


def _build_system_prompt(context: "ExecutionContext", active_doc: Optional[dict] = None) -> str:
    schema = get_allowed_endpoints_schema()

    user_ctx = context.academic_context or {}
    ctx_parts = []
    if context.user_id:
        ctx_parts.append(f"userId={context.user_id}")
    if context.role:
        ctx_parts.append(f"role={context.role}")
    for key in ("batchId", "departmentId", "collegeId", "studentId", "batchCode", "fullName", "studentName"):
        if val := user_ctx.get(key):
            ctx_parts.append(f"{key}={val}")

    ctx_line = " | ".join(ctx_parts)

    role_persona = _load_role_persona(context.role or "")
    persona_section = f"\n## الشخصية والأسلوب (مطلوب الالتزام به)\n{role_persona}\n" if role_persona else ""

    # ── Memory/profile/entity context ────────────────────────────────────
    meta = context.metadata or {}
    user_ctx = context.academic_context or {}
    memory_section = ""
    profile_section = ""
    entity_section = ""
    goal_section = ""

    memory = meta.get("memory") or {}
    if memory:
        parts = []
        if memory.get("last_intent"):
            parts.append(f"- آخر طلب: {memory['last_intent']}")
        if memory.get("last_message"):
            parts.append(f"- آخر رسالة: {str(memory['last_message'])[:120]}")
        if memory.get("last_result"):
            parts.append(f"- آخر نتيجة (ملخص): {str(memory['last_result'])[:250]}")
        if parts:
            memory_section = "## ذاكرة المحادثة السابقة\n" + "\n".join(parts) + "\n"

    academic_profile = meta.get("academic_profile") or {}
    personalized = meta.get("personalized_context") or ""
    if personalized:
        profile_section = f"## الملف الأكاديمي للمستخدم\n{personalized}\n"
    elif academic_profile:
        profile_section = f"## الملف الأكاديمي\n{academic_profile}\n"

    # Conversation entities — course names, doctor names, goals mentioned earlier
    conv_entities = meta.get("conversation_entities") or user_ctx.get("conversation_entities") or {}
    if conv_entities:
        try:
            from app.core.entity_tracker import build_entity_context_block
            block = build_entity_context_block(conv_entities)
            if block:
                entity_section = block + "\n"
        except Exception:
            pass

    # User goal tracking
    user_goal = meta.get("user_goal") or user_ctx.get("user_goal") or ""
    if user_goal:
        goal_labels = {
            "graduation": "التخرج قريباً",
            "improve_gpa": "تحسين الـ GPA",
            "exam_prep": "التحضير للامتحانات",
            "understand_topic": "فهم موضوع محدد",
            "registration": "التسجيل في المواد",
        }
        goal_label = goal_labels.get(user_goal, user_goal)
        goal_section = f"## هدف المستخدم الحالي\n- {goal_label}\n"

    active_doc_section = _build_active_doc_section(active_doc or {})

    return f"""أنت **مساعد أكاديمي ذكي** — المساعد الشخصي للطالب في كل شيء يخص جامعته.
أنت تعمل كـ reasoning agent: تفكر، تستخدم أدواتك، وتحل أي مشكلة للطالب.

{persona_section}
## بيانات الجلسة الحالية
{ctx_line}
{memory_section}{profile_section}{entity_section}{goal_section}{active_doc_section}

## هويتك وقدراتك الكاملة
أنت تستطيع المساعدة في كل ما يحتاجه الطالب:
- **بياناته الأكاديمية**: درجاته، جدوله، مواده المسجلة، GPA
- **محتوى المحاضرات**: قراءة ملفات PDF/PPT وشرحها وتلخيصها
- **التعلم النشط**: شرح مفاهيم، إجابة أسئلة، توليد أمثلة، مقارنة أفكار
- **اللائحة الأكاديمية**: الخطة الدراسية، متطلبات التخرج، ساعات الـ credit
- **التواصل**: إرسال شكاوى أو رسائل للدكاترة والإدارة
- **المتابعة**: حالة الطلبات، الردود، الإشعارات

## كيف تتصرف — الفلسفة الأساسية

**أنت لا تستسلم أبداً.** كل طلب له حل:
- لو API فشلت → جرب endpoint تاني
- لو الملف مش قابل للقراءة → اشرح الموضوع من معرفتك العامة (ووضح إنك بتشرح من معرفتك مش من الملف)
- لو مش عارف المقصود → اسأل سؤال ذكي واحد يوضح الصورة
- لو البيانات فارغة → قل إنها فارغة واقترح البديل

**أنت تستخدم أدواتك بذكاء:**
- للبيانات من الجامعة: `call_backend_api`
- لقراءة ملف PDF: `read_material_pdf` (بعد ما تجيب الـ file_url من API)
- لاللائحة الدراسية: `read_regulation_pdf`
- لأسئلة عامة عن مفاهيم دراسية: أجب مباشرة من معرفتك

## تسلسل التفكير
1. **افهم** ماذا يريد الطالب بالضبط
2. **قرر** هل تحتاج بيانات من الجامعة أم تجيب من معرفتك
3. **استخدم الأدوات** إذا احتجت بيانات — وبالتوازي إذا كانت متعددة
4. **أجب** بشكل مفيد ومنظم

## حل الضمائر والسياق
لو الرسالة قصيرة أو فيها ضمير ("ها"، "فيها"، "it"):
- ابحث في المحادثة السابقة لتحديد المرجع
- "شرحها" بعد ذكر "قواعد البيانات" = اشرح قواعد البيانات
- "كمّل" = استمر في نفس الموضوع
- لا تسأل لو الجواب واضح من السياق

## للمواد والملفات — الطريقة الصحيحة
لو الطالب طلب شرح مادة أو ملف:
1. جيب قائمة مواده: `GET /api/Enrollments/my-enrollments` → يرجع `subjectCode`
2. جيب الملفات: `GET /api/Materials/by-subject/{{subjectCode}}`
3. اقرأ الملف: `read_material_pdf` مع الـ `fileUrl`
4. اشرح بناءً على المحتوى

⚡ "مواد دكتور حمدي" = مادة DS-101 → جيب ملفاتها واقرأها واشرح
⚡ "اشرح أول لكتشر في الـ CV" = جيب CV-00 materials → اقرأ → اشرح

## ممنوعات بيانات الجامعة
- لا تخترع أرقاماً (درجات، إحصائيات) لم تأتِ من API
- إذا فارغة → قل إنها فارغة
- إذا خطأ API → جرب بديل، لا تستسلم

## اللغة
تكلم بنفس لغة المستخدم بالضبط — عربي بعربي، إنجليزي بإنجليزي.

## نقاط النهاية المتاحة
{schema}"""


# ── Follow-up chain detection ─────────────────────────────────────────────────

def _is_followup_message(message: str) -> bool:
    """Return True if the message is a short follow-up command to act on a document."""
    msg = message.strip().lower()
    if any(kw in msg for kw in _FOLLOWUP_AR):
        return True
    if any(kw in msg for kw in _FOLLOWUP_EN):
        return True
    return False


def _inject_followup_context(message: str, active_doc: Optional[dict]) -> str:
    """
    If the message is a follow-up command AND there is an active material document,
    prepend a context hint so the LLM uses read_material_pdf with the stored URL.
    Returns the (possibly modified) message.
    """
    if not active_doc:
        return message
    if active_doc.get("document_type") != "material":
        return message
    file_url = active_doc.get("file_url", "").strip()
    if not file_url:
        return message
    if _is_followup_message(message):
        prefix = f"[CONTEXT: Active file = {file_url}. Use read_material_pdf with this URL.] "
        logger.info(
            "[ReactAgent] follow-up detected — injecting file_url into message url=%.60s",
            file_url,
        )
        return prefix + message
    return message


# ── ReactAgent ────────────────────────────────────────────────────────────────

class ReactAgent:
    """
    ReAct-style agent using gpt-4o-mini native function calling.

    run()        → returns str  (used by the normal /api/chat endpoint)
    stream_run() → yields str tokens (used by /api/chat/stream endpoint)
    """

    def __init__(self, openrouter_client: Any, model_router: Any) -> None:
        self.client = openrouter_client
        self.model_router = model_router

    # ── Normal (non-streaming) path ───────────────────────────────────────────

    async def run(self, context: "ExecutionContext") -> str:
        """Execute the ReAct loop and return the final response string."""
        context.set_model(_MODEL)
        context.set_tool("react_agent")

        # Normalize academic_context keys (snake_case ↔ camelCase)
        if context.academic_context:
            context.academic_context = _normalize_academic_context(context.academic_context)

        # Input guard — truncate suspiciously long messages / injection attempts
        input_check = check_user_input(context.message)
        if not input_check.passed:
            logger.warning(
                "[ReactAgent] input guard: risk=%s warnings=%s user_id=%s",
                input_check.risk_level, input_check.warnings, context.user_id,
            )
            if input_check.risk_level == "high":
                return "لم نتمكن من معالجة هذا الطلب. يرجى إعادة الصياغة."
            # Use sanitized (truncated) version
            context.message = input_check.sanitized_text

        # Fast-path: answer identity questions without any LLM call
        fast = _try_answer_from_context(context)
        if fast:
            logger.info("[ReactAgent] fast-path answer user_id=%s", context.user_id)
            return fast

        # Load active document context from memory store
        active_doc: Optional[dict] = None
        if context.user_id:
            try:
                from app.services.memory_store import get_memory_store
                _mem = get_memory_store()
                active_doc = await _mem.get_active_document(context.user_id)
                if active_doc:
                    logger.info(
                        "[ReactAgent] loaded active_doc for user_id=%s type=%s title=%s",
                        context.user_id, active_doc.get("document_type"), active_doc.get("title"),
                    )
                    # Also inject into metadata so _tool_read_material fallback can use it
                    if not (context.metadata or {}).get("last_file_url") and active_doc.get("file_url"):
                        context.add_metadata("last_file_url", active_doc["file_url"])
            except Exception as _exc:
                logger.warning("[ReactAgent] failed to load active_doc: %s", _exc)

        # Follow-up chain detection (Task 4):
        # If the user message is a short follow-up command AND there is an active
        # material document, prepend context so the LLM uses read_material_pdf.
        context.message = _inject_followup_context(context.message, active_doc)

        messages = _build_messages(context, active_doc=active_doc)
        self._collected_results: list = []   # accumulate tool outputs for guard

        logger.info(
            "[ReactAgent] START user_id=%s role=%s message=%.100r",
            context.user_id, context.role, context.message,
        )
        t_start = time.perf_counter()
        _tool_call_count = 0

        for iteration in range(1, _MAX_ITERATIONS + 1):
            try:
                response = await self.client.chat.completions.create(
                    model=_MODEL,
                    messages=messages,
                    tools=_ALL_TOOLS,
                    tool_choice="auto",
                )
            except Exception as exc:
                logger.error("[ReactAgent] LLM error iteration=%d — %s", iteration, exc)
                return "حدث خطأ أثناء المعالجة. حاول مرة أخرى."

            choice = response.choices[0]
            finish = choice.finish_reason

            logger.debug("[ReactAgent] iteration=%d finish_reason=%s", iteration, finish)

            # ── Final answer ──────────────────────────────────────────────
            if finish in ("stop", "end_turn") or (
                finish is None and not getattr(choice.message, "tool_calls", None)
            ):
                answer = choice.message.content or ""

                # Self-reflection: if answer is too short + tools were available,
                # push one more iteration asking model to elaborate or fetch data.
                # Only fires if we still have budget AND haven't already used tools.
                if (
                    len(answer.strip()) < 60
                    and _tool_call_count == 0
                    and iteration < _MAX_ITERATIONS
                ):
                    logger.info(
                        "[ReactAgent] short answer without tool calls (len=%d) — "
                        "triggering self-reflection iteration=%d",
                        len(answer), iteration,
                    )
                    messages.append({"role": "assistant", "content": answer})
                    messages.append({
                        "role": "user",
                        "content": (
                            "[تعليمات داخلية — لا تعرضها للمستخدم]: "
                            "إجابتك قصيرة جداً. إذا كانت البيانات مطلوبة، "
                            "استدعِ tool مناسب الآن. إذا لم تكن بيانات مطلوبة، "
                            "وسّع الإجابة بشكل أكثر فائدة وتفصيلاً."
                        ),
                    })
                    continue  # retry with self-reflection

                elapsed = round(time.perf_counter() - t_start, 3)

                # Response guard — log hallucination risks without blocking
                guard = guard_response(answer, self._collected_results, context.message)
                if not guard.passed:
                    logger.warning(
                        "[ReactAgent] response guard: risk=%s warnings=%s user_id=%s",
                        guard.risk_level, guard.warnings, context.user_id,
                    )

                logger.info(
                    "[ReactAgent] END iterations=%d tool_calls=%d duration_s=%s "
                    "len=%d guard=%s user_id=%s",
                    iteration, _tool_call_count, elapsed,
                    len(answer), guard.risk_level, context.user_id,
                )
                # Store metrics on context for upstream observability
                context.add_metadata("react_iterations", iteration)
                context.add_metadata("react_tool_calls", _tool_call_count)
                context.add_metadata("react_guard_risk", guard.risk_level)
                return answer

            # ── Parallel tool calls ───────────────────────────────────────
            if getattr(choice.message, "tool_calls", None):
                messages.append(_serialize_msg(choice.message))

                tool_calls = choice.message.tool_calls
                _tool_call_count += len(tool_calls)
                logger.info(
                    "[ReactAgent] iteration=%d dispatching %d tool(s) in parallel",
                    iteration, len(tool_calls),
                )

                results = await asyncio.gather(
                    *[_dispatch_tool(tc, context, self.model_router) for tc in tool_calls],
                    return_exceptions=True,
                )

                for tc, result in zip(tool_calls, results):
                    if isinstance(result, Exception):
                        logger.warning("[ReactAgent] tool %s raised: %s", tc.function.name, result)
                        result = {"error": str(result), "tool": tc.function.name}
                    # Collect for response guard
                    self._collected_results.append(result)
                    result_str = _format_tool_result(tc.function.name, result)
                    messages.append({
                        "role": "tool",
                        "tool_call_id": tc.id,
                        "content": result_str,
                    })

        logger.warning("[ReactAgent] max iterations (%d) reached", _MAX_ITERATIONS)
        return "تعذّر إتمام الطلب بعد عدة محاولات. جرّب صياغة مختلفة."

    # ── Streaming path ────────────────────────────────────────────────────────

    async def stream_run(self, context: "ExecutionContext") -> AsyncGenerator[str, None]:
        """
        Streaming version of run().

        Tool-calling iterations execute normally (no streaming — tool calls
        require complete JSON responses).  The FINAL answer is streamed
        token-by-token via OpenAI streaming so the user sees the first word
        within ~500ms of the last tool call completing.
        """
        # Normalize academic_context keys (snake_case ↔ camelCase)
        if context.academic_context:
            context.academic_context = _normalize_academic_context(context.academic_context)

        context.set_model(_MODEL)
        context.set_tool("react_agent")

        # Fast-path: yield instantly without any LLM call
        fast = _try_answer_from_context(context)
        if fast:
            yield fast
            return

        # Load active document context
        active_doc: Optional[dict] = None
        if context.user_id:
            try:
                from app.services.memory_store import get_memory_store
                _mem = get_memory_store()
                active_doc = await _mem.get_active_document(context.user_id)
                if active_doc and not (context.metadata or {}).get("last_file_url") and active_doc.get("file_url"):
                    context.add_metadata("last_file_url", active_doc["file_url"])
            except Exception as _exc:
                logger.warning("[ReactAgent.stream] failed to load active_doc: %s", _exc)

        # Follow-up chain detection
        context.message = _inject_followup_context(context.message, active_doc)

        messages = _build_messages(context, active_doc=active_doc)

        logger.info(
            "[ReactAgent.stream] START user_id=%s message=%.80r",
            context.user_id, context.message,
        )

        for iteration in range(1, _MAX_ITERATIONS + 1):
            is_last = iteration == _MAX_ITERATIONS

            if is_last:
                # Force a direct text answer on the last iteration (no tools)
                # and stream it token-by-token.
                try:
                    stream = await self.client.chat.completions.create(
                        model=_MODEL,
                        messages=messages,
                        stream=True,
                        # No tools param → model cannot call tools, must answer directly
                    )
                    collected: list[str] = []
                    async for chunk in stream:
                        content = chunk.choices[0].delta.content
                        if content:
                            collected.append(content)
                            yield content
                    logger.info(
                        "[ReactAgent.stream] END (last iter streamed) len=%d",
                        sum(len(c) for c in collected),
                    )
                except Exception as exc:
                    logger.error("[ReactAgent.stream] final stream error: %s", exc)
                    yield "حدث خطأ. حاول مرة أخرى."
                return

            # Non-final iteration: use normal (non-streaming) call with tools
            try:
                response = await self.client.chat.completions.create(
                    model=_MODEL,
                    messages=messages,
                    tools=_ALL_TOOLS,
                    tool_choice="auto",
                )
            except Exception as exc:
                logger.error("[ReactAgent.stream] LLM error it=%d: %s", iteration, exc)
                yield "حدث خطأ. حاول مرة أخرى."
                return

            choice = response.choices[0]
            finish = choice.finish_reason

            # Final answer reached before max iterations — stream it directly
            if finish in ("stop", "end_turn") or not getattr(choice.message, "tool_calls", None):
                answer = choice.message.content or ""
                logger.info(
                    "[ReactAgent.stream] END (early finish it=%d) len=%d",
                    iteration, len(answer),
                )
                # Stream token-by-token via the model for true live output
                try:
                    messages.append({"role": "assistant", "content": answer})
                    # Re-ask with streaming so the user sees live tokens;
                    # we send the already-computed answer as assistant turn so the
                    # model just needs to "continue" (in practice it echoes it
                    # streaming — same content, live delivery).
                    # Simpler and zero-hallucination: replay the known text in
                    # small chunks at natural typing speed.
                    words = answer.split()
                    for i in range(0, len(words), 3):
                        chunk = " ".join(words[i:i + 3]) + " "
                        yield chunk
                        await asyncio.sleep(0.015)
                except Exception as exc:
                    logger.warning("[ReactAgent.stream] replay error: %s", exc)
                    yield answer
                return

            # ── Parallel tool calls ───────────────────────────────────────
            messages.append(_serialize_msg(choice.message))
            tool_calls = choice.message.tool_calls
            results = await asyncio.gather(
                *[_dispatch_tool(tc, context, self.model_router) for tc in tool_calls],
                return_exceptions=True,
            )
            for tc, result in zip(tool_calls, results):
                if isinstance(result, Exception):
                    result = {"error": str(result), "tool": tc.function.name}
                result_str = _format_tool_result(tc.function.name, result)
                messages.append({
                    "role": "tool",
                    "tool_call_id": tc.id,
                    "content": result_str,
                })

        yield "تعذّر إتمام الطلب بعد عدة محاولات."


# ── Tool result formatter ─────────────────────────────────────────────────────

def _format_tool_result(tool_name: str, result: Any) -> str:
    """
    Format a tool result into a structured observation string for the LLM.

    Raw JSON blobs are hard for the model to reason about. This function
    produces a short, labelled observation that includes the key data
    while staying within token budget.
    """
    if isinstance(result, dict) and result.get("error"):
        return f"[{tool_name}] ❌ خطأ: {result['error']}"

    # Academic analysis: return the full analysis text directly — it's already formatted
    if isinstance(result, dict) and result.get("academic_analysis"):
        return f"[{tool_name}] ✅ التحليل الأكاديمي:\n{result['academic_analysis']}"

    # Regulation answer: return directly
    if isinstance(result, dict) and result.get("regulation_answer"):
        return f"[{tool_name}] ✅ من اللائحة:\n{result['regulation_answer']}"

    try:
        raw = json.dumps(result, ensure_ascii=False, default=str)
    except Exception:
        raw = str(result)

    # Truncate very long results, keeping the beginning (most important data)
    if len(raw) > 5_000:
        raw = raw[:5_000] + "\n... [truncated]"

    # For API responses, try to extract a useful summary line
    if isinstance(result, dict):
        data = result.get("data") or result
        total = None
        if isinstance(data, dict):
            total = data.get("total") or data.get("count")
        elif isinstance(data, list):
            total = len(data)

        if total is not None:
            return f"[{tool_name}] ✅ النتيجة: {total} سجل\n{raw}"

    return f"[{tool_name}] ✅\n{raw}"


# ── Message builder ───────────────────────────────────────────────────────────

def _build_messages(context: "ExecutionContext", active_doc: Optional[dict] = None) -> List[Dict[str, Any]]:
    msgs: List[Dict[str, Any]] = [
        {"role": "system", "content": _build_system_prompt(context, active_doc=active_doc)}
    ]

    # Use last 10 turns (increased from 8 for better multi-turn reasoning)
    history_window = (context.history or [])[-10:]
    for turn in history_window:
        role = turn.get("role", "")
        content = str(turn.get("content", "")).strip()
        if role in ("user", "assistant") and content:
            msgs.append({"role": role, "content": content[:1000]})
    msgs.append({"role": "user", "content": context.message})
    return msgs


# ── Tool dispatcher ───────────────────────────────────────────────────────────

async def _dispatch_tool(tool_call: Any, context: "ExecutionContext", model_router: Any) -> Any:
    fn = tool_call.function.name
    try:
        args = json.loads(tool_call.function.arguments or "{}")
    except json.JSONDecodeError:
        return {"error": "invalid JSON in tool arguments"}

    logger.info("[ReactAgent] tool=%s args=%.200s", fn, str(args))

    if fn == "call_backend_api":
        return await _tool_call_api(args, context)
    elif fn == "read_regulation_pdf":
        return await _tool_regulation(args, context, model_router)
    elif fn == "read_material_pdf":
        return await _tool_read_material(args, context, model_router)
    elif fn == "generate_exam":
        return await _tool_generate_exam(args, context, model_router)
    elif fn == "analyze_academic_profile":
        return await _tool_academic_analysis(args, context, model_router)
    else:
        return {"error": f"unknown tool: {fn}"}


# ── Material context extractor ────────────────────────────────────────────────

async def _extract_and_save_material_context(
    result: Any, context: "ExecutionContext"
) -> None:
    """
    Scan an API result for fileUrl / file_url items and save to active_document.

    Called after call_backend_api succeeds. Handles the common shape:
      {"data": {"items": [{"fileUrl": "...", "title": "...", "subjectName": "..."}, ...]}}
    """
    try:
        if not isinstance(result, dict):
            return
        data = result.get("data") or result
        items: list = []
        if isinstance(data, dict):
            items = data.get("items") or data.get("results") or []
            if not items and isinstance(data, dict):
                # Sometimes the response IS the item dict directly
                items = [data]
        elif isinstance(data, list):
            items = data

        for item in items:
            if not isinstance(item, dict):
                continue
            file_url = (item.get("fileUrl") or item.get("file_url") or "").strip()
            if not file_url:
                continue
            title = (
                item.get("title") or item.get("fileName") or item.get("name")
                or file_url.split("/")[-1].split("?")[0]
            )
            subject_name = (
                item.get("subjectName") or item.get("subject_name")
                or item.get("subject") or ""
            )
            material_id = str(item.get("id") or item.get("materialId") or "")
            active_doc = {
                "document_type": "material",
                "file_url": file_url,
                "title": title,
                "material_id": material_id,
                "subject_name": subject_name,
            }
            from app.services.memory_store import get_memory_store
            _mem = get_memory_store()
            await _mem.set_active_document(context.user_id or "", active_doc)
            # Also update the last_file metadata so the fallback in _tool_read_material works
            context.add_metadata("last_file_url", file_url)
            logger.info(
                "[ReactAgent] auto-saved active_doc from API result: title=%s url=%.60s",
                title, file_url,
            )
            break  # save the first found; subsequent items are secondary
    except Exception as exc:
        logger.warning("[ReactAgent] _extract_and_save_material_context failed: %s", exc)


# ── Tool: call_backend_api ────────────────────────────────────────────────────

async def _tool_call_api(args: dict, context: "ExecutionContext") -> Any:
    from app.services.backend_client import tool_execution_client as client

    method = (args.get("method") or "GET").upper()
    path   = (args.get("path") or "").strip()
    qp     = args.get("query_params") or {}
    body   = args.get("body") or {}
    auth   = (context.metadata or {}).get("auth_header") or ""

    if not path:
        return {"error": "path is required"}

    if not validate_endpoint(method, path):
        logger.warning("[ReactAgent] BLOCKED %s %s", method, path)
        return {"error": f"Endpoint {method} {path} is not in the allowed list."}

    try:
        if method == "GET":
            result = await client.fetch(route=path, auth_header=auth, params=qp or None)
        elif method == "POST":
            result = await client.post(route=path, payload=body, auth_header=auth)
        else:
            return {"error": f"Method {method} not supported"}

        logger.info("[ReactAgent] API %s %s → ok", method, path)
        # Auto-save any material file URL found in the response
        if context.user_id:
            asyncio.create_task(
                _extract_and_save_material_context(result, context)
            )
        return result

    except Exception as exc:
        logger.warning("[ReactAgent] API %s %s failed — %s", method, path, exc)
        return {"error": str(exc), "path": path}


# ── Tool: read_regulation_pdf ─────────────────────────────────────────────────

async def _tool_regulation(args: dict, context: "ExecutionContext", model_router: Any) -> Any:
    try:
        from app.agents.schemas import AgentInput
        from app.modules.regulation import RegulationModule

        auth = (context.metadata or {}).get("auth_header") or ""
        module = RegulationModule(
            model_router=model_router,
            backend_client=__import__(
                "app.services.backend_client", fromlist=["tool_execution_client"]
            ).tool_execution_client,
        )
        agent_input = AgentInput(
            user_id=context.user_id or "",
            message=args.get("query") or context.message,
            auth_header=auth,
            context={
                "selected_model": _MODEL,
                "history": context.history or [],
            },
        )
        output = await module.run(agent_input)
        return {"regulation_answer": output.response, "status": output.status}

    except Exception as exc:
        logger.error("[ReactAgent] read_regulation_pdf failed — %s", exc)
        return {"error": str(exc)}


# ── Tool: read_material_pdf ───────────────────────────────────────────────────

async def _tool_read_material(args: dict, context: "ExecutionContext", model_router: Any) -> Any:
    """Download and summarize/read a lecture material PDF."""
    import httpx
    import io

    task     = args.get("task", "summarize")
    file_url = (args.get("file_url") or "").strip()
    question = args.get("question") or ""
    auth     = (context.metadata or {}).get("auth_header") or ""

    # Fallback: check agent memory for last saved file URL
    if not file_url:
        mem = context.metadata or {}
        file_url = (mem.get("last_file_url") or "").strip()

    if not file_url:
        return {
            "error": (
                "file_url is required but was not provided. "
                "First call the backend API to get the file URL, then call read_material_pdf."
            )
        }

    # 2. Download PDF
    try:
        async with httpx.AsyncClient(timeout=30.0) as http:
            resp = await http.get(file_url)
            resp.raise_for_status()
            raw = resp.content
    except Exception as exc:
        logger.error("[read_material_pdf] download failed url=%s — %s", file_url, exc)
        return {"error": f"Could not download the file: {exc}"}

    # 3. Extract text — try multiple methods
    fname = file_url.split("?")[0].split("/")[-1].lower()
    text = ""

    # Method A: pdfminer (best quality for text PDFs)
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract
        text = pdfminer_extract(io.BytesIO(raw)) or ""
    except Exception:
        pass

    # Method B: pypdf fallback
    if not text.strip():
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(raw))
            text = "\n\n".join(p.extract_text() or "" for p in reader.pages)
        except Exception:
            pass

    # Method C: PPT/PPTX text extraction
    if not text.strip() and fname.endswith((".ppt", ".pptx")):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw))
            slides = []
            for i, slide in enumerate(prs.slides):
                slide_text = " ".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text")
                )
                if slide_text.strip():
                    slides.append(f"[Slide {i+1}] {slide_text}")
            text = "\n".join(slides)
        except Exception:
            pass

    # If text extracted is very short → use deep PDF mode for large PDFs
    if len(text.strip()) < 200 and fname.endswith(".pdf") and len(raw) > 50_000:
        try:
            from app.modules.deep_pdf_study import extract_pdf_pages, summarize_large_pdf
            pages = extract_pdf_pages(raw)
            if pages:
                q = question or task or "اشرح محتوى الملف"
                result = await summarize_large_pdf(
                    pages=pages, user_question=q,
                    model_router=model_router, model_id=_MODEL,
                )
                return {"material_content": result, "chars_processed": sum(len(p) for p in pages)}
        except Exception as dp_exc:
            logger.warning("[read_material_pdf] deep_pdf failed — %s", dp_exc)

    # If still empty → try vision for scanned PDFs / images
    if not text.strip():
        try:
            from app.modules.file_extraction import _vision_extract
            vision_text = await _vision_extract(raw, fname, question or "اشرح محتوى الملف", model_router, auth)
            if vision_text and vision_text.strip():
                text = vision_text
        except Exception:
            pass

    # If completely empty → inform AI to use general knowledge
    if not text.strip():
        fname_display = file_url.split("?")[0].split("/")[-1]
        return {
            "error": "unreadable_file",
            "message": (
                f"الملف '{fname_display}' مش قابل للقراءة كـ نص (ممكن يكون PDF مسحوب ضوئياً أو slides بصور). "
                "اشرح للطالب الموضوع من معرفتك الأكاديمية العامة وأخبره إنك مش بتشرح من الملف."
            )
        }

    # Cap text at 20K chars for better coverage
    text = text[:20_000]

    # 4. Build task-specific prompt
    focus = f"\nالطالب سأل عن: {question}" if question else ""
    if task == "list_headings":
        user_prompt = f"اعرض الموضوعات والعناوين الرئيسية في هذا الملف:\n{focus}\n\n{text}"
    elif task == "explain":
        user_prompt = f"اشرح محتوى هذا الملف بالتفصيل بطريقة مبسطة للطالب:{focus}\n\n{text}"
    elif task == "read":
        user_prompt = f"اقرأ هذا الملف وقدم المعلومات الأساسية منه:{focus}\n\n{text}"
    else:  # summarize / default
        user_prompt = f"قدم شرحاً وافياً ومفصلاً لهذا الملف الأكاديمي:{focus}\n\n{text}"

    # 5. LLM narration
    try:
        result = await model_router.generate_with_messages(
            messages=[
                {"role": "system", "content": (
                    "أنت مدرّس أكاديمي متخصص. اشرح محتوى الملف بعمق وتفصيل. "
                    "استخدم عناوين وأمثلة واشرح كل مفهوم بوضوح. "
                    "أجب بنفس لغة الطالب (عربي أو إنجليزي)."
                )},
                {"role": "user", "content": user_prompt},
            ],
            model_id=_MODEL,
            max_tokens=2000,
        )
        return {"material_content": result, "chars_processed": len(text)}
    except Exception as exc:
        logger.error("[read_material_pdf] LLM failed — %s", exc)
        return {"error": str(exc)}


# ── Tool: generate_exam ───────────────────────────────────────────────────────

async def _tool_generate_exam(args: dict, context: "ExecutionContext", model_router: Any) -> Any:
    if context.role not in ("doctor", "admin", "superadmin"):
        return {"error": "Permission denied: only doctors and admins can generate exams."}

    try:
        from app.agents.schemas import AgentInput, ExecutionPlan, ExamParams
        from app.modules.exam_generation import ExamGenerationModule

        auth = (context.metadata or {}).get("auth_header") or ""

        exam_params = ExamParams(
            subject=args.get("subject", ""),
            question_count=args.get("question_count", 10),
            difficulty=args.get("difficulty", "medium"),
            exam_type=args.get("exam_type", "mcq"),
            topics=args.get("topics") or [],
        )
        if soid := args.get("subject_offering_id"):
            context.academic_context["subjectOfferingId"] = soid

        plan = ExecutionPlan(
            goal_summary="Generate exam via ReactAgent.",
            intent="generate_exam",
            is_executable=True,
            exam_params=exam_params,
        )
        module = ExamGenerationModule(
            model_router=model_router,
            backend_client=__import__(
                "app.services.backend_client", fromlist=["tool_execution_client"]
            ).tool_execution_client,
        )
        agent_input = AgentInput(
            user_id=context.user_id or "",
            message=context.message,
            auth_header=auth,
            context={
                "selected_model": "openai/gpt-4o",
                "academic_context": context.academic_context,
                "history": context.history or [],
            },
        )
        output = await module.run(agent_input, plan=plan)
        return {"exam_result": output.response, "status": output.status}

    except Exception as exc:
        logger.error("[ReactAgent] generate_exam failed — %s", exc)
        return {"error": str(exc)}


# ── Tool: analyze_academic_profile ───────────────────────────────────────────

async def _tool_academic_analysis(
    args: dict, context: "ExecutionContext", model_router: Any
) -> Any:
    """
    Deep academic analysis tool — calls AcademicAdvisorModule internally.

    Fetches: roadmap + student overview + regulation RAG passages.
    Returns a structured analysis string the ReactAgent presents to the user.
    """
    if context.role not in ("student",):
        # Admins/doctors use different analysis flows
        return {"error": "analyze_academic_profile is only available for students."}

    try:
        from app.agents.schemas import AgentInput
        from app.modules.academic_advisor import AcademicAdvisorModule

        auth    = (context.metadata or {}).get("auth_header") or ""
        focus   = args.get("focus", "general")
        question = args.get("question") or context.message

        # Build a focused question that combines focus + original question
        focus_prompts = {
            "graduation_readiness": "هل أنا مؤهل للتخرج؟ ما الشروط المتبقية؟",
            "gpa_improvement":      "كيف أحسّن معدلي الأكاديمي؟ ما المواد الحرجة؟",
            "failed_subjects":      "ما المواد اللي رسبت فيها ومحتاج أعيدها؟",
            "next_semester_plan":   "ما المواد المقترحة للترم الجاي؟",
            "academic_risk":        "هل أنا في خطر أكاديمي؟ ما حدود الإنذار؟",
            "general":              question,
        }
        focused_question = focus_prompts.get(focus, question)
        if question and question != context.message:
            focused_question = f"{question} ({focused_question})"

        module = AcademicAdvisorModule(
            model_router=model_router,
            backend_client=__import__(
                "app.services.backend_client", fromlist=["tool_execution_client"]
            ).tool_execution_client,
        )
        agent_input = AgentInput(
            user_id=context.user_id or "",
            message=focused_question,
            auth_header=auth,
            context={
                "selected_model": "openai/gpt-4o-mini",
                "academic_context": context.academic_context or {},
                "history": context.history or [],
            },
        )
        output = await module.run(agent_input)
        if output.status == "success":
            return {
                "academic_analysis": output.response,
                "focus": focus,
                "data": output.data or {},
            }
        return {"error": output.response or "Academic analysis failed."}

    except Exception as exc:
        logger.error("[ReactAgent] analyze_academic_profile failed — %s", exc, exc_info=True)
        return {"error": str(exc)}


# ── Helper ────────────────────────────────────────────────────────────────────

def _serialize_msg(msg: Any) -> dict:
    """Convert OpenAI message object → plain dict for the messages list."""
    d: dict = {"role": msg.role}
    if msg.content:
        d["content"] = msg.content
    if tcs := getattr(msg, "tool_calls", None):
        d["tool_calls"] = [
            {
                "id": tc.id,
                "type": "function",
                "function": {
                    "name": tc.function.name,
                    "arguments": tc.function.arguments,
                },
            }
            for tc in tcs
        ]
    return d
